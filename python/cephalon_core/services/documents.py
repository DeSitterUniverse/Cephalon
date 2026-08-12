import csv
from dataclasses import dataclass, field
import hashlib
import os
import time
import uuid

import docx
import openpyxl
import pptx
from .. import storage
from ..validators import is_supported_file
from .pdf_parser import DocumentBlock, PdfAsset, PARSER_VERSION as PDF_PARSER_VERSION, parse_pdf
from .table_models import MAX_TABLE_ROWS, StructuredTable, TABLE_PARSER_VERSION, build_table

SKIPPED_DIRECTORY_NAMES = {
    ".git",
    ".obsidian",
    ".trash",
    ".venv",
    "__pycache__",
    "node_modules",
}


@dataclass
class ExtractedDocument:
    text: str
    extraction_mode: str
    blocks: list[DocumentBlock] = field(default_factory=list)
    page_count: int | None = None
    warnings: list[str] = field(default_factory=list)
    parser_version: str = "cephalon-basic-2026-05"
    assets: list[PdfAsset] = field(default_factory=list)
    tables: list[StructuredTable] = field(default_factory=list)


def get_file_hash(path: str) -> str:
    hasher = hashlib.sha256()
    with open(path, "rb") as f:
        while chunk := f.read(1024 * 1024):
            hasher.update(chunk)
    return hasher.hexdigest()


def file_metadata(path: str) -> tuple[int, int]:
    stat = os.stat(path)
    return stat.st_size, int(stat.st_mtime)


def find_existing_doc_by_hash(sqlite_conn, content_hash: str):
    return storage.fetchone(
        sqlite_conn,
        "SELECT id, path, status, chunk_count FROM documents WHERE content_hash = ? AND type = 'file' AND status IN ('ready', 'ingesting') LIMIT 1",
        (content_hash,),
    )


def looks_like_text(path: str, sample_size: int = 8192) -> bool:
    with open(path, "rb") as f:
        sample = f.read(sample_size)
    if not sample:
        return True
    if b"\x00" in sample:
        return False
    control_bytes = sum(1 for byte in sample if byte < 32 and byte not in {9, 10, 13})
    return control_bytes / len(sample) < 0.08


def read_text_fallback(path: str) -> str:
    for encoding in ("utf-8-sig", "utf-16", "utf-8", "latin-1"):
        try:
            with open(path, "r", encoding=encoding) as f:
                return f.read()
        except UnicodeError:
            continue
    with open(path, "r", encoding="latin-1", errors="replace") as f:
        return f.read()


def extract_text(path: str, force_text: bool = False) -> tuple[str, str]:
    ext = os.path.splitext(path)[1].lower()

    if force_text and not looks_like_text(path):
        raise ValueError("File appears to be binary and cannot be safely imported as text.")
    if force_text:
        return read_text_fallback(path), "text"

    if ext == ".pdf":
        return parse_pdf(path).text, "native_structured"
    if ext == ".docx":
        doc = docx.Document(path)
        return "\n".join([para.text for para in doc.paragraphs]), "native"
    if ext == ".pptx":
        prs = pptx.Presentation(path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs), "native"
    if ext == ".xlsx":
        wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
        text_runs = []
        for sheet in wb.worksheets:
            text_runs.append(f"--- Sheet: {sheet.title} ---")
            for row in sheet.iter_rows(values_only=True):
                row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                if row_text.strip():
                    text_runs.append(row_text)
        return "\n".join(text_runs), "native"
    if ext == ".csv":
        text_runs = []
        with open(path, newline="", encoding="utf-8") as f:
            reader = csv.reader(f)
            for row in reader:
                text_runs.append("\t".join(row))
        return "\n".join(text_runs), "native"

    if not looks_like_text(path):
        raise ValueError("Unknown file type appears to be binary and cannot be safely imported as text.")
    return read_text_fallback(path), "text"


def extract_document(path: str, force_text: bool = False) -> ExtractedDocument:
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf" and not force_text:
        parsed = parse_pdf(path)
        return ExtractedDocument(
            text=parsed.text,
            extraction_mode="native_structured",
            blocks=parsed.blocks,
            page_count=parsed.page_count,
            warnings=parsed.warnings,
            parser_version=parsed.parser_version,
            assets=parsed.assets,
            tables=parsed.tables,
        )

    if ext == ".csv" and not force_text:
        return _extract_csv_document(path)
    if ext == ".xlsx" and not force_text:
        return _extract_xlsx_document(path)

    text, extraction_mode = extract_text(path, force_text=force_text)
    block = DocumentBlock(text=text, page_number=1) if text.strip() else None
    return ExtractedDocument(
        text=text,
        extraction_mode=extraction_mode,
        blocks=[block] if block else [],
        page_count=1 if block else 0,
        parser_version=PDF_PARSER_VERSION if ext == ".pdf" else "cephalon-basic-2026-05",
    )


def _extract_csv_document(path: str) -> ExtractedDocument:
    warnings: list[str] = []
    handle = None
    for encoding in ("utf-8-sig", "utf-8", "utf-16", "latin-1"):
        try:
            handle = open(path, newline="", encoding=encoding)
            sample = handle.read(8192)
            handle.seek(0)
            break
        except UnicodeError:
            if handle is not None:
                handle.close()
            handle = None
    if handle is None:
        raise ValueError("CSV encoding could not be decoded safely.")
    if encoding not in {"utf-8-sig", "utf-8"}:
        warnings.append(f"csv_encoding_fallback:{encoding}")
    with handle:
        try:
            dialect = csv.Sniffer().sniff(sample, delimiters=",\t;|")
        except csv.Error:
            dialect = csv.excel
            warnings.append("csv_dialect_fallback")
        rows = []
        for index, row in enumerate(csv.reader(handle, dialect)):
            if index >= MAX_TABLE_ROWS:
                warnings.append("row_limit_reached")
                break
            rows.append(row)
    table = build_table(
        rows,
        source_type="csv",
        table_index=0,
        provenance={"encoding": encoding, "delimiter": dialect.delimiter},
        warnings=warnings,
    )
    block = DocumentBlock(
        text=table.text,
        page_number=1,
        block_type="table",
        provenance={"table_index": 0, "source_type": "csv"},
        structured_table=table,
    )
    return ExtractedDocument(
        text=table.text,
        extraction_mode="native_structured",
        blocks=[block] if table.text.strip() else [],
        page_count=1,
        warnings=warnings,
        parser_version=TABLE_PARSER_VERSION,
        tables=[table],
    )


def _extract_xlsx_document(path: str) -> ExtractedDocument:
    workbook = openpyxl.load_workbook(path, data_only=False, read_only=False)
    warnings: list[str] = []
    tables: list[StructuredTable] = []
    blocks: list[DocumentBlock] = []
    try:
        for sheet_index, sheet in enumerate(workbook.worksheets):
            rows: list[list[object]] = []
            formulas: dict[tuple[int, int], str] = {}
            number_formats: dict[tuple[int, int], str] = {}
            for row_index, cells in enumerate(sheet.iter_rows()):
                if row_index >= MAX_TABLE_ROWS:
                    warnings.append(f"sheet:{sheet.title}:row_limit_reached")
                    break
                values = []
                for column_index, cell in enumerate(cells):
                    value = cell.value
                    if cell.data_type == "f" and value is not None:
                        formula = str(value)
                        formulas[(row_index, column_index)] = formula if formula.startswith("=") else f"={formula}"
                    number_formats[(row_index, column_index)] = str(cell.number_format or "")
                    values.append(value)
                rows.append(values)
            while rows and not any(value not in (None, "") for value in rows[-1]):
                rows.pop()
            if not rows:
                continue
            merged_ranges = [str(cell_range) for cell_range in sheet.merged_cells.ranges]
            table = build_table(
                rows,
                source_type="xlsx",
                table_index=sheet_index,
                sheet_name=sheet.title,
                sheet_index=sheet_index,
                formulas=formulas,
                number_formats=number_formats,
                merged_ranges=merged_ranges,
                provenance={"formula_mode": "preserve_formula_no_recalculation"},
            )
            tables.append(table)
            blocks.append(DocumentBlock(
                text=f"--- Sheet: {sheet.title} ---\n{table.text}",
                page_number=sheet_index + 1,
                block_type="table",
                provenance={"table_index": sheet_index, "source_type": "xlsx", "sheet_name": sheet.title},
                structured_table=table,
            ))
    finally:
        workbook.close()
    return ExtractedDocument(
        text="\n\n".join(block.text for block in blocks),
        extraction_mode="native_structured",
        blocks=blocks,
        page_count=len(blocks),
        warnings=warnings,
        parser_version=TABLE_PARSER_VERSION,
        tables=tables,
    )


def collect_supported_files(path: str, force_text: bool = False) -> list[str]:
    if os.path.isfile(path):
        return [path]

    files: list[str] = []
    for root, dirs, names in os.walk(path):
        dirs[:] = [name for name in dirs if name not in SKIPPED_DIRECTORY_NAMES]
        for name in names:
            full_path = os.path.join(root, name)
            if is_supported_file(full_path) or looks_like_text(full_path) or force_text:
                files.append(full_path)
    return sorted(files)


def collect_obsidian_files(vault_path: str) -> list[str]:
    if not os.path.isdir(vault_path):
        return []

    files: list[str] = []
    for root, dirs, names in os.walk(vault_path):
        dirs[:] = [
            name for name in dirs
            if name not in SKIPPED_DIRECTORY_NAMES and not name.startswith(".")
        ]
        for name in names:
            if name.startswith("."):
                continue
            full_path = os.path.join(root, name)
            ext = os.path.splitext(name)[1].lower()
            if ext in {".md", ".txt", ".canvas"} or looks_like_text(full_path):
                files.append(full_path)
    return sorted(files)


def register_ingesting_document(sqlite_conn, path: str, content_hash: str, extraction_mode: str = "native", doc_id: str | None = None, embedding_metadata: dict | None = None) -> str:
    doc_id = doc_id or str(uuid.uuid4())
    embedding_metadata = embedding_metadata or storage.active_embedding_metadata()
    size_bytes, modified_at = file_metadata(path)
    storage.execute(
        sqlite_conn,
        """
        INSERT INTO documents
            (id, path, display_name, content_hash, ingested_at, chunk_count, status, type, size_bytes, modified_at, last_indexed_at, extraction_mode, embedding_model_id, embedding_dim, stale_embedding)
        VALUES (?, ?, COALESCE((SELECT display_name FROM documents WHERE id = ?), ?), ?, ?, ?, ?, 'file', ?, ?, ?, ?, ?, ?, 0)
        ON CONFLICT(id) DO UPDATE SET
            path = excluded.path,
            content_hash = excluded.content_hash,
            ingested_at = excluded.ingested_at,
            chunk_count = excluded.chunk_count,
            status = excluded.status,
            size_bytes = excluded.size_bytes,
            modified_at = excluded.modified_at,
            last_indexed_at = excluded.last_indexed_at,
            extraction_mode = excluded.extraction_mode,
            embedding_model_id = excluded.embedding_model_id,
            embedding_dim = excluded.embedding_dim,
            stale_embedding = 0,
            stale_reasons = NULL,
            last_error = NULL
        """,
        (
            doc_id,
            path,
            doc_id,
            os.path.basename(path),
            content_hash,
            int(time.time()),
            0,
            "ingesting",
            size_bytes,
            modified_at,
            int(time.time()),
            extraction_mode,
            embedding_metadata["embedding_model_id"],
            embedding_metadata["embedding_dim"],
        ),
    )
    return doc_id


def mark_document_ready(sqlite_conn, doc_id: str, chunk_count: int) -> None:
    storage.execute(
        sqlite_conn,
        "UPDATE documents SET status = 'ready', chunk_count = ?, stale_embedding = 0, stale_reasons = NULL, last_error = NULL, last_indexed_at = ? WHERE id = ?",
        (chunk_count, int(time.time()), doc_id),
    )


def mark_document_failed(sqlite_conn, doc_id: str, error: str) -> None:
    storage.execute(
        sqlite_conn,
        "UPDATE documents SET status = ?, chunk_count = 0, last_error = ? WHERE id = ?",
        (f"failed: {error[:50]}", error, doc_id),
    )
